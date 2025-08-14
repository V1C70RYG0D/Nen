import re
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def read_text_file(path: Path) -> str:
	return path.read_text(encoding="utf-8")


def extract_first_group(pattern: re.Pattern, text: str) -> Optional[str]:
	match = pattern.search(text)
	return match.group(1).strip() if match else None


def parse_slides_from_markdown(markdown_text: str) -> List[Dict[str, Any]]:
	lines = markdown_text.splitlines()
	slides: List[Dict[str, Any]] = []

	# Identify slide section headers like: #### Slide 1: Title – Subtitle
	header_regex = re.compile(r"^####\s*Slide\s*(\d+):\s*(.+?)\s*$")
	idx = 0
	while idx < len(lines):
		line = lines[idx].rstrip("\n")
		header_match = header_regex.match(line)
		if not header_match:
			idx += 1
			continue

		# Capture current slide header and advance
		slide_number = header_match.group(1)
		slide_title_line = header_match.group(2)

		# Collect block until next slide header or EOF
		block_start = idx + 1
		idx += 1
		while idx < len(lines) and not header_regex.match(lines[idx]):
			idx += 1
		block_lines = lines[block_start:idx]

		# Parse content bullets under "Content (Layer 2)"
		content_header_regex = re.compile(r"^\s*-\s*\*\*Content \(Layer 2\)\*\*:\s*$")
		presenter_regex = re.compile(r"^\s*-\s*\*\*Presenter Script \(Layer 2\)\*\*:\s*(.+?)\s*$")
		purpose_regex = re.compile(r"^\s*-\s*\*\*Purpose \(Layer 1\)\*\*:\s*(.+?)\s*$")
		hook_regex = re.compile(r"^\s*-\s*\*\*Teaser Hook \(Layer 1\)\*\*:\s*(.+?)\s*$")
		speaker_note_regex = re.compile(r"^\s*-\s*\*\*Speaker Note \(Layer 3\)\*\*:\s*(.+?)\s*$")

		bullets: List[str] = []
		presenter_script: Optional[str] = None
		purpose: Optional[str] = None
		teaser_hook: Optional[str] = None
		speaker_note: Optional[str] = None

		in_content = False
		for raw in block_lines:
			if content_header_regex.match(raw):
				in_content = True
				continue
			# Nested bullets are indented by at least 2 spaces followed by '- '
			if in_content:
				if re.match(r"^\s{2,}-\s+", raw):
					text = re.sub(r"^\s*-\s+", "", raw).strip()
					# Keep label + description if present (e.g., **Logo**: ...)
					bullets.append(text)
					continue
				# End of content section when next top-level bullet appears
				if re.match(r"^\s*-\s+\*\*.+\*\*:.*$", raw):
					in_content = False
					# fall through to evaluate the same line for other matches

			# Capture single-line metadata
			if presenter_script is None:
				m = presenter_regex.match(raw)
				if m:
					presenter_script = m.group(1).strip()
					continue
			if purpose is None:
				m = purpose_regex.match(raw)
				if m:
					purpose = m.group(1).strip()
					continue
			if teaser_hook is None:
				m = hook_regex.match(raw)
				if m:
					teaser_hook = m.group(1).strip()
					continue
			if speaker_note is None:
				m = speaker_note_regex.match(raw)
				if m:
					speaker_note = m.group(1).strip()
					continue

		# Build the slide title
		slide_title = slide_title_line
		slides.append(
			{
				"number": int(slide_number),
				"title": slide_title,
				"bullets": bullets,
				"notes": build_notes(presenter_script, purpose, teaser_hook, speaker_note),
			}
		)

	# Sort by slide number to ensure correct order
	slides.sort(key=lambda s: s["number"]) 
	return slides


def build_notes(presenter: Optional[str], purpose: Optional[str], hook: Optional[str], speaker_note: Optional[str]) -> str:
	parts: List[str] = []
	if presenter:
		parts.append(f"Presenter Script: {presenter}")
	if purpose:
		parts.append(f"Purpose: {purpose}")
	if hook:
		parts.append(f"Teaser Hook: {hook}")
	if speaker_note:
		parts.append(f"Speaker Note: {speaker_note}")
	return "\n".join(parts)


def apply_dark_theme(slide):
	# Black background
	fill = slide.background.fill
	fill.solid()
	fill.fore_color.rgb = RGBColor(0, 0, 0)


def set_title_text(title_shape, text: str):
	if not title_shape:
		return
	title_shape.text = text
	p = title_shape.text_frame.paragraphs[0]
	p.alignment = PP_ALIGN.LEFT
	font = p.font
	font.bold = True
	font.size = Pt(40)
	font.color.rgb = RGBColor(255, 255, 255)


def set_bullets(placeholder, bullets: List[str]):
	if not placeholder:
		return
	tf = placeholder.text_frame
	# Clear any existing text by assigning the first bullet to the first paragraph
	if bullets:
		tf.clear() if hasattr(tf, "clear") else None
		# Ensure at least one paragraph exists
		if len(tf.paragraphs) == 0:
			_ = tf.add_paragraph()
		# Add bullets
		for idx, bullet in enumerate(bullets):
			if idx == 0:
				p = tf.paragraphs[0]
				p.text = bullet
			else:
				p = tf.add_paragraph()
				p.text = bullet
			p.level = 0
			pf = p.font
			pf.size = Pt(22)
			pf.color.rgb = RGBColor(255, 255, 255)


def set_notes(slide, notes_text: Optional[str]):
	if not notes_text:
		return
	notes_slide = slide.notes_slide
	tf = notes_slide.notes_text_frame
	tf.text = notes_text


def create_presentation(slides_data: List[Dict[str, Any]], output_path: Path):
	prs = Presentation()
	# Set 16:9 aspect ratio
	prs.slide_width = Inches(13.333)
	prs.slide_height = Inches(7.5)

	for s in slides_data:
		# Use Title and Content layout (index 1)
		slide_layout = prs.slide_layouts[1]
		slide = prs.slides.add_slide(slide_layout)
		apply_dark_theme(slide)

		# Title
		title_shape = slide.shapes.title
		set_title_text(title_shape, s.get("title", f"Slide {s.get('number', '')}"))

		# Bullets
		content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
		set_bullets(content_placeholder, s.get("bullets", []))

		# Notes
		set_notes(slide, s.get("notes"))

	prs.save(str(output_path))


def main():
	root = Path(__file__).parent
	md_path = Path(sys.argv[1]) if len(sys.argv) > 1 else root / "Pitch-deck.md"
	out_path = Path(sys.argv[2]) if len(sys.argv) > 2 else root / "pitch_deck.pptx"

	if not md_path.exists():
		raise FileNotFoundError(f"Markdown file not found: {md_path}")

	md_text = read_text_file(md_path)
	slides = parse_slides_from_markdown(md_text)
	if not slides:
		raise RuntimeError("No slides parsed from markdown. Ensure headings start with '#### Slide N: ...'")

	create_presentation(slides, out_path)
	print(f"Generated: {out_path}")


if __name__ == "__main__":
	main()


